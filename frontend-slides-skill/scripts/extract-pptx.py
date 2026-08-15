#!/usr/bin/env python3
"""
extract-pptx.py — pull text, speaker notes, and images out of an existing
.pptx into a structured outline, so a frontend-slides deck can be rebuilt
from it as HTML instead of PowerPoint.

Requires: pip install python-pptx

Usage:
    python3 extract-pptx.py deck.pptx --out extracted/
    python3 extract-pptx.py deck.pptx --out extracted/ --json

Output:
    extracted/outline.md        Human-readable outline, one section per slide.
    extracted/outline.json      Same content as structured JSON (with --json).
    extracted/images/           Every picture in the deck, named
                                 slide<N>-img<M>.<ext>.
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print(
        "error: python-pptx is required. Install it with:\n"
        "    pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(1)


def shape_text(shape):
    if not shape.has_text_frame:
        return None
    text = "\n".join(
        p.text for p in shape.text_frame.paragraphs if p.text.strip()
    )
    return text.strip() or None


def is_title_shape(shape, slide):
    return slide.shapes.title is not None and shape.shape_id == slide.shapes.title.shape_id


def extract_table(shape):
    rows = []
    for row in shape.table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return rows


def extract_slide(slide, index, images_dir):
    title = None
    body_texts = []
    tables = []
    images = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            ext = image.ext or "png"
            filename = f"slide{index}-img{len(images) + 1}.{ext}"
            (images_dir / filename).write_bytes(image.blob)
            images.append(filename)
            continue

        if shape.has_table:
            tables.append(extract_table(shape))
            continue

        text = shape_text(shape)
        if not text:
            continue

        if title is None and is_title_shape(shape, slide):
            title = text
        else:
            body_texts.append(text)

    notes = None
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        notes = notes_text or None

    return {
        "index": index,
        "title": title,
        "body": body_texts,
        "tables": tables,
        "images": images,
        "notes": notes,
    }


def render_markdown(slides):
    lines = ["# Extracted outline", ""]
    for s in slides:
        lines.append(f"## Slide {s['index']}" + (f" — {s['title']}" if s["title"] else ""))
        lines.append("")
        for text in s["body"]:
            for line in text.splitlines():
                lines.append(f"- {line}")
        for table in s["tables"]:
            for row in table:
                lines.append("| " + " | ".join(row) + " |")
            lines.append("")
        if s["images"]:
            lines.append("")
            lines.append("Images: " + ", ".join(f"`images/{img}`" for img in s["images"]))
        if s["notes"]:
            lines.append("")
            lines.append(f"> Speaker notes: {s['notes']}")
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def main():
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("pptx", type=Path, help="Path to the source .pptx file")
    parser.add_argument("--out", type=Path, default=Path("extracted"), help="Output directory (default: extracted/)")
    parser.add_argument("--json", action="store_true", help="Also write outline.json")
    args = parser.parse_args()

    if not args.pptx.exists():
        print(f"error: {args.pptx} not found", file=sys.stderr)
        sys.exit(1)

    images_dir = args.out / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    presentation = Presentation(args.pptx)
    slides = [
        extract_slide(slide, i, images_dir)
        for i, slide in enumerate(presentation.slides, start=1)
    ]

    outline_md = args.out / "outline.md"
    outline_md.write_text(render_markdown(slides), encoding="utf-8")
    print(f"wrote {outline_md} ({len(slides)} slides)")

    if args.json:
        outline_json = args.out / "outline.json"
        outline_json.write_text(json.dumps(slides, indent=2), encoding="utf-8")
        print(f"wrote {outline_json}")

    image_count = sum(len(s["images"]) for s in slides)
    if image_count:
        print(f"extracted {image_count} image(s) to {images_dir}")


if __name__ == "__main__":
    main()
